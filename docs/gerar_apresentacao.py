# -*- coding: utf-8 -*-
"""
Gera a apresentação final do projeto Tech Hub (Marketplace de Tecnologia).
Segue a estrutura sugerida no Guia de Projeto (Etapa 8 - Apresentação Final).

Uso:
    .venv/Scripts/python.exe docs/gerar_apresentacao.py

Saída:
    docs/Tech_Hub_Apresentacao.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# PALETA — extraída de static/css/style.css (tema escuro)
# ============================================================
PRIMARY_DARK   = RGBColor(0x04, 0x11, 0x23)  # fundo
SECONDARY_DARK = RGBColor(0x08, 0x1D, 0x3D)  # cards
ACCENT_CYAN    = RGBColor(0x3A, 0xF6, 0xFF)  # destaque
PRIMARY_BLUE   = RGBColor(0x0F, 0xA7, 0xFF)
TEXT_WHITE     = RGBColor(0xEE, 0xF7, 0xFF)
TEXT_GRAY      = RGBColor(0x9B, 0xB1, 0xCA)
TEXT_MUTED     = RGBColor(0x6B, 0x7A, 0x99)
SUCCESS        = RGBColor(0x00, 0xFF, 0x88)
WARNING        = RGBColor(0xFF, 0xC1, 0x07)
CARD_BORDER    = RGBColor(0x1E, 0x3A, 0x5F)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ============================================================
# Helpers
# ============================================================
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=PRIMARY_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None, line_w=None, radius=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return tb, tf


def add_run(p, text, size, color, bold=False, italic=False, font="Segoe UI"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def para(tf, first=False):
    if first and not tf.paragraphs[0].runs:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def accent_bar(slide):
    """Faixa fina ciano→azul no topo (igual ao header do site)."""
    bar = rect(slide, 0, 0, SLIDE_W, Inches(0.10), ACCENT_CYAN)
    _grad_horizontal(bar, ACCENT_CYAN, PRIMARY_BLUE)


def _grad_horizontal(shape, c1, c2):
    """Aplica gradiente horizontal simples num shape via XML."""
    spPr = shape.fill._xPr
    # remove solidFill existente
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn('a:gradFill'), {})
    gsLst = grad.makeelement(qn('a:gsLst'), {})
    for pos, color in ((0, c1), (100000, c2)):
        gs = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        srgb = grad.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
        gs.append(srgb)
        gsLst.append(gs)
    grad.append(gsLst)
    lin = grad.makeelement(qn('a:lin'), {'ang': '0', 'scaled': '1'})
    grad.append(lin)
    # inserir após a:prstGeom/ no início do spPr (antes de a:ln)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


def section_header(slide, etapa, titulo, meta=None, accent=PRIMARY_BLUE):
    """Cabeçalho de seção estilo o PDF do guia."""
    band = rect(slide, Inches(0.55), Inches(0.75), Inches(12.23), Inches(1.15),
                SECONDARY_DARK, radius=True)
    # barra colorida à esquerda
    rect(slide, Inches(0.55), Inches(0.75), Inches(0.14), Inches(1.15), accent, radius=False)
    tb, tf = textbox(slide, Inches(0.95), Inches(0.82), Inches(9.5), Inches(1.0),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    add_run(p, etapa, 11, accent, bold=True)
    p2 = para(tf)
    add_run(p2, titulo, 25, TEXT_WHITE, bold=True)
    if meta:
        tbm, tfm = textbox(slide, Inches(9.7), Inches(0.82), Inches(2.9), Inches(1.0),
                           anchor=MSO_ANCHOR.MIDDLE)
        tfm.paragraphs[0].alignment = PP_ALIGN.RIGHT
        add_run(tfm.paragraphs[0], meta, 11, TEXT_GRAY)


def bullets(slide, x, y, w, h, items, size=15, gap=10, anchor=MSO_ANCHOR.TOP):
    tb, tf = textbox(slide, x, y, w, h, anchor=anchor)
    for i, item in enumerate(items):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if isinstance(item, tuple):
            label, desc = item
            add_run(p, "▸  ", size, ACCENT_CYAN, bold=True)
            add_run(p, label, size, TEXT_WHITE, bold=True)
            if desc:
                add_run(p, "  —  " + desc, size, TEXT_GRAY)
        else:
            add_run(p, "▸  ", size, ACCENT_CYAN, bold=True)
            add_run(p, item, size, TEXT_GRAY)
    return tb


def card(slide, x, y, w, h, title, lines, accent=ACCENT_CYAN, title_size=15, body_size=12):
    c = rect(slide, x, y, w, h, SECONDARY_DARK, line=CARD_BORDER, line_w=Pt(1), radius=True)
    rect(slide, x, y, w, Inches(0.09), accent, radius=False)
    tb, tf = textbox(slide, x + Inches(0.22), y + Inches(0.20),
                     w - Inches(0.44), h - Inches(0.36))
    p = para(tf, first=True)
    p.space_after = Pt(7)
    add_run(p, title, title_size, TEXT_WHITE, bold=True)
    for ln in lines:
        pp = para(tf)
        pp.space_after = Pt(4)
        pp.line_spacing = 1.04
        if isinstance(ln, tuple):
            add_run(pp, ln[0] + "  ", body_size, accent, bold=True)
            add_run(pp, ln[1], body_size, TEXT_GRAY)
        else:
            add_run(pp, "• ", body_size, accent)
            add_run(pp, ln, body_size, TEXT_GRAY)
    return c


def footer(slide, idx):
    tb, tf = textbox(slide, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    add_run(p, "Tech Hub", 9, ACCENT_CYAN, bold=True)
    add_run(p, "  ·  Marketplace de Tecnologia  ·  Programação Back-End", 9, TEXT_MUTED)
    tb2, tf2 = textbox(slide, Inches(11.8), Inches(7.05), Inches(1.0), Inches(0.35),
                       anchor=MSO_ANCHOR.MIDDLE)
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    add_run(tf2.paragraphs[0], str(idx), 9, TEXT_MUTED)


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
s = add_slide()
set_bg(s)
# blocos decorativos
rect(s, Inches(9.7), Inches(-1.2), Inches(5), Inches(5), SECONDARY_DARK)
big = rect(s, Inches(10.4), Inches(2.1), Inches(2.4), Inches(2.4), PRIMARY_DARK,
           line=ACCENT_CYAN, line_w=Pt(1.5), radius=True)
tbq, tfq = textbox(s, Inches(10.4), Inches(2.1), Inches(2.4), Inches(2.4),
                   anchor=MSO_ANCHOR.MIDDLE)
tfq.paragraphs[0].alignment = PP_ALIGN.CENTER
add_run(tfq.paragraphs[0], "🛒", 54, ACCENT_CYAN)
accent_bar(s)

tb, tf = textbox(s, Inches(0.8), Inches(2.2), Inches(8.8), Inches(3.2))
p = para(tf, first=True)
add_run(p, "TECH HUB", 54, TEXT_WHITE, bold=True)
p2 = para(tf)
p2.space_before = Pt(4)
add_run(p2, "Marketplace de Tecnologia", 26, ACCENT_CYAN, bold=True)
p3 = para(tf)
p3.space_before = Pt(14)
add_run(p3, "Compra, venda e troca de hardware entre pessoas e lojas —\ncom negociação, pagamento e entrega de ponta a ponta.",
        15, TEXT_GRAY)
p4 = para(tf)
p4.space_before = Pt(22)
add_run(p4, "Programação Back-End  ·  Django  ·  Apresentação Final", 13, TEXT_MUTED)

# ============================================================
# SLIDE 2 — O QUE É O SISTEMA / PROBLEMA
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "1 · O QUE É O SISTEMA", "O problema que resolve", meta="Visão geral", accent=PRIMARY_BLUE)
tb, tf = textbox(s, Inches(0.6), Inches(2.15), Inches(12.1), Inches(0.9))
p = para(tf, first=True)
add_run(p, "Comprar e trocar produtos de tecnologia usados é confuso e inseguro: ", 16, TEXT_WHITE)
add_run(p, "anúncios espalhados, negociação solta no chat e sem garantia de pagamento ou entrega.",
        16, TEXT_GRAY)
card(s, Inches(0.6), Inches(3.25), Inches(3.9), Inches(3.4), "O Tech Hub centraliza",
     ["Anúncios de venda E de troca no mesmo lugar",
      "Negociação estruturada por propostas",
      "Pagamento e entrega acompanhados",
      "Perfis de pessoa física e loja"], accent=ACCENT_CYAN)
card(s, Inches(4.72), Inches(3.25), Inches(3.9), Inches(3.4), "Para quem é",
     [("PF", "compra, vende e troca usados/seminovos"),
      ("Loja", "vende produtos novos com selo verificado"),
      ("Comprador", "negocia e paga com segurança"),
      ("Vendedor", "controla anúncios e entregas")], accent=PRIMARY_BLUE)
card(s, Inches(8.84), Inches(3.25), Inches(3.9), Inches(3.4), "Diferencial",
     ["Sistema completo de negociação de trocas",
      "Propostas e contrapropostas por turnos",
      "Confirmação de entrega dos 2 lados",
      "Pagamento de diferença em dinheiro"], accent=SUCCESS)
footer(s, 2)

# ============================================================
# SLIDE 3 — PERFIS E REGRAS DE NEGÓCIO
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "REGRAS DE NEGÓCIO", "Perfis de usuário e permissões", meta="RF01–RF03", accent=ACCENT_CYAN)
card(s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(4.4), "👤  Pessoa Física (CPF)",
     ["Compra de lojas e de outros usuários",
      "Cria anúncios de Venda, Troca ou ambos",
      "Pode anunciar produtos usados/seminovos",
      "Participa de negociações de troca",
      "Cadastro validado: CPF único + maior de 18 anos"],
     accent=ACCENT_CYAN, title_size=18, body_size=13)
card(s, Inches(6.82), Inches(2.2), Inches(5.9), Inches(4.4), "🏢  Loja (CNPJ)",
     ["Vende produtos (apenas novos)",
      "Cria anúncios somente de Venda",
      "Não participa de trocas",
      "Selo visual de \"Loja Verificada\"",
      "Cadastro validado: CNPJ único + dados do responsável"],
     accent=PRIMARY_BLUE, title_size=18, body_size=13)
footer(s, 3)

# ============================================================
# SLIDE 4 — ARQUITETURA E STACK
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "DECISÕES TÉCNICAS", "Arquitetura e stack", meta="RNF01–RNF02", accent=WARNING)
card(s, Inches(0.6), Inches(2.2), Inches(3.85), Inches(2.05), "Backend",
     [("Python", "linguagem"), ("Django 5.2", "framework web"),
      ("DRF + SimpleJWT", "API REST e auth")], accent=ACCENT_CYAN)
card(s, Inches(4.66), Inches(2.2), Inches(3.85), Inches(2.05), "Dados",
     [("PostgreSQL", "banco relacional"), ("Pillow", "upload de imagens"),
      ("django-environ", "config via .env")], accent=PRIMARY_BLUE)
card(s, Inches(8.72), Inches(2.2), Inches(3.85), Inches(2.05), "Frontend",
     [("Templates Django", "server-side"), ("CSS próprio", "tema claro/escuro"),
      ("JS puro", "máscaras e validação")], accent=SUCCESS)
card(s, Inches(0.6), Inches(4.45), Inches(5.95), Inches(2.2),
     "Organização por domínios",
     ["Views finas em marketplace_app/domains/:",
      ("auth", "cadastro, login, perfil"),
      ("listings", "anúncios e busca"),
      ("cart / checkout", "carrinho, pedidos, pagamento"),
      ("trades", "negociação de trocas")],
     accent=WARNING, body_size=12)
card(s, Inches(6.78), Inches(4.45), Inches(5.95), Inches(2.2),
     "Integrações externas",
     [("Mercado Pago", "checkout + webhook de pagamento"),
      ("API REST", "listagem pública e cadastro via JSON"),
      ("JWT", "autenticação stateless para a API"),
      ("GitHub Actions", "CI: migrations + testes + lint")],
     accent=ACCENT_CYAN, body_size=12)
footer(s, 4)

# ============================================================
# SLIDE 5 — MODELAGEM DE DADOS
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "MODELAGEM DO BANCO", "Principais entidades", meta="PostgreSQL", accent=RGBColor(0x8A,0x4F,0xD4))
card(s, Inches(0.6), Inches(2.2), Inches(3.9), Inches(4.4), "Identidade",
     [("User", "AbstractUser + is_store"),
      ("CommonProfile", "CPF, nascimento, endereço"),
      ("StoreProfile", "CNPJ, razão social, selo"),
      ("Category", "categorias de hardware")],
     accent=ACCENT_CYAN, body_size=13)
card(s, Inches(4.72), Inches(2.2), Inches(3.9), Inches(4.4), "Catálogo & compra",
     [("Listing", "anúncio: venda/troca"),
      ("ListingImage", "múltiplas imagens"),
      ("Comment", "comentários aninhados"),
      ("Cart / CartItem", "carrinho"),
      ("Order / OrderItem", "pedido + snapshots"),
      ("Payment / Delivery", "pagamento e entrega")],
     accent=PRIMARY_BLUE, body_size=13)
card(s, Inches(8.84), Inches(2.2), Inches(3.9), Inches(4.4), "Negociação de troca",
     [("TradeRequest", "negociação"),
      ("TradeProposal", "proposta + dinheiro"),
      ("TradeProposalImage", "fotos da proposta"),
      ("TradeMessage", "mensagens"),
      ("TradeFulfillment", "execução/acordo"),
      ("TradeDelivery", "entrega dos 2 lados")],
     accent=SUCCESS, body_size=13)
footer(s, 5)

# ============================================================
# SLIDE 6 — FUNCIONALIDADES: COMPRA
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "2 · DEMONSTRAÇÃO", "Fluxo de compra", meta="RF02–RF04", accent=PRIMARY_BLUE)
steps = [
    ("1", "Buscar", "Busca + filtros por categoria, condição, preço e tipo de vendedor"),
    ("2", "Anúncio", "Detalhes, galeria de imagens e comentários com respostas"),
    ("3", "Carrinho", "Adiciona itens marcando compra ou troca"),
    ("4", "Checkout", "Endereço de entrega + pagamento (Mercado Pago)"),
    ("5", "Pedido", "Acompanha status do pagamento e da entrega"),
]
x = Inches(0.6)
cw = Inches(2.34)
gap = Inches(0.13)
for num, title, desc in steps:
    c = rect(s, x, Inches(2.5), cw, Inches(3.6), SECONDARY_DARK, line=CARD_BORDER, line_w=Pt(1), radius=True)
    circ = slide_circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + cw/2 - Inches(0.4), Inches(2.2), Inches(0.8), Inches(0.8))
    circ.fill.solid(); circ.fill.fore_color.rgb = PRIMARY_BLUE; _no_line(circ); circ.shadow.inherit = False
    ctf = circ.text_frame; ctf.word_wrap = True
    add_run(ctf.paragraphs[0], num, 22, TEXT_WHITE, bold=True)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tb, tf = textbox(s, x + Inches(0.12), Inches(2.95), cw - Inches(0.24), Inches(3.0))
    p = para(tf, first=True); p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8)
    add_run(p, title, 16, ACCENT_CYAN, bold=True)
    p2 = para(tf); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.1
    add_run(p2, desc, 12, TEXT_GRAY)
    x = Emu(int(x) + int(cw) + int(gap))
footer(s, 6)

# ============================================================
# SLIDE 7 — DIFERENCIAL: NEGOCIAÇÃO DE TROCA
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "2 · DEMONSTRAÇÃO  ·  DIFERENCIAL", "Sistema de negociação de troca",
               meta="RF05", accent=SUCCESS)
tb, tf = textbox(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(0.7))
p = para(tf, first=True)
add_run(p, "Negociação por turnos: ", 15, TEXT_WHITE, bold=True)
add_run(p, "os participantes alternam propostas e contrapropostas até um acordo — "
           "com itens, imagens e diferença em dinheiro.", 15, TEXT_GRAY)
card(s, Inches(0.6), Inches(2.95), Inches(5.95), Inches(3.7), "Como funciona",
     [("1.", "Comprador abre a negociação no anúncio de troca"),
      ("2.", "Envia proposta: item + valor em dinheiro + fotos"),
      ("3.", "O outro lado aceita ou faz contraproposta"),
      ("4.", "Qualquer um dos dois pode aceitar a proposta atual"),
      ("5.", "Aceito → checkout e confirmação de entrega"),
      ("6.", "Os 2 confirmam a entrega → troca concluída")],
     accent=SUCCESS, title_size=16, body_size=13)
card(s, Inches(6.78), Inches(2.95), Inches(5.95), Inches(3.7), "Estados da negociação",
     [("Aguardando", "proposta inicial pendente"),
      ("Em negociação", "propostas indo e voltando"),
      ("Aprovada", "acordo fechado, rumo à entrega"),
      ("Concluída", "entrega confirmada pelos dois"),
      ("Recusada / Cancelada", "encerrada sem acordo")],
     accent=ACCENT_CYAN, title_size=16, body_size=13)
footer(s, 7)

# ============================================================
# SLIDE 8 — QUALIDADE, VALIDAÇÕES E SEGURANÇA
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "QUALIDADE & SEGURANÇA", "Validações, testes e proteção", meta="RNF + Etapa 6", accent=RGBColor(0xE0,0x4F,0x5F))
card(s, Inches(0.6), Inches(2.2), Inches(3.9), Inches(4.4), "Validações de regra",
     ["CPF e CNPJ com dígito verificador",
      "Documentos únicos no cadastro",
      "Maioridade (18+) na pessoa física",
      "Loja não cria troca / só produto novo",
      "Não adicionar o próprio anúncio ao carrinho",
      "Máscaras de CPF, telefone, CEP e R$"],
     accent=ACCENT_CYAN, body_size=12.5)
card(s, Inches(4.72), Inches(2.2), Inches(3.9), Inches(4.4), "Testes automatizados",
     ["Cadastro de PF e PJ",
      "Fluxo de checkout e pedidos",
      "Fluxo completo de troca",
      "Entrega e histórico",
      "Webhooks de pagamento",
      "Auth e regras dos models"],
     accent=SUCCESS, body_size=12.5)
card(s, Inches(8.84), Inches(2.2), Inches(3.9), Inches(4.4), "Segurança",
     ["Senhas validadas pelo Django",
      "Segredos fora do código (.env)",
      "CSRF em todos os formulários",
      "Cookies seguros + SSL em produção",
      "HSTS e ALLOWED_HOSTS obrigatórios",
      "JWT para a API REST"],
     accent=WARNING, body_size=12.5)
footer(s, 8)

# ============================================================
# SLIDE 9 — DECISÕES TÉCNICAS (POR QUE ASSIM)
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "3 · DECISÕES TÉCNICAS", "Por que fizemos assim", meta="2–3 min", accent=PRIMARY_BLUE)
bullets(s, Inches(0.7), Inches(2.25), Inches(12.0), Inches(4.5), [
    ("User customizado desde o início", "AbstractUser com is_store evita migração dolorosa depois e unifica PF/PJ num só login."),
    ("Perfis separados (Common/Store)", "cada tipo tem campos e validações próprios sem poluir o User com dezenas de colunas opcionais."),
    ("Lógica em domains/ e não nas views", "views finas que delegam para módulos de domínio — mais fácil de ler, testar e dividir entre o grupo."),
    ("Snapshots em OrderItem", "título e preço são copiados no pedido; o histórico não muda se o anúncio for editado depois."),
    ("Troca como entidades próprias", "TradeRequest/Proposal/Fulfillment modelam a negociação real, em vez de forçar tudo no chat."),
    ("Pagamento via gateway + webhook", "o Mercado Pago confirma o pagamento de forma assíncrona, refletindo no status do pedido."),
], size=14.5, gap=11)
footer(s, 9)

# ============================================================
# SLIDE 10 — DIFICULDADES E APRENDIZADOS
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "4 · DIFICULDADES", "O que foi difícil e o que aprendemos", meta="1–2 min", accent=WARNING)
card(s, Inches(0.6), Inches(2.2), Inches(5.95), Inches(4.4), "Desafios",
     ["Modelar a negociação de troca por turnos",
      "Garantir que os dois lados confirmem a entrega",
      "Integrar pagamento e tratar o webhook",
      "Validar CPF/CNPJ e aplicar máscaras",
      "Manter o visual consistente em todas as telas"],
     accent=RGBColor(0xE0,0x4F,0x5F), title_size=18, body_size=13.5)
card(s, Inches(6.78), Inches(2.2), Inches(5.95), Inches(4.4), "Aprendizados",
     ["Modelagem de dados orienta todo o resto",
      "Separar a lógica em domínios facilita o trabalho em grupo",
      "Testes dão segurança para mexer no código",
      "Validar no backend, não só no frontend",
      "Git com commits frequentes salva o projeto"],
     accent=SUCCESS, title_size=18, body_size=13.5)
footer(s, 10)

# ============================================================
# SLIDE 11 — PRÓXIMOS PASSOS
# ============================================================
s = add_slide(); set_bg(s); accent_bar(s)
section_header(s, "ROADMAP", "Próximos passos", meta="Etapa 7 · Deploy", accent=ACCENT_CYAN)
card(s, Inches(0.6), Inches(2.2), Inches(3.9), Inches(4.4), "Em andamento",
     ["Deploy em produção (Railway / Render)",
      "collectstatic + Postgres gerenciado",
      "Variáveis de ambiente no servidor",
      "Domínio e HTTPS"],
     accent=WARNING, body_size=13)
card(s, Inches(4.72), Inches(2.2), Inches(3.9), Inches(4.4), "Melhorias planejadas",
     ["Avaliações de vendedor (RF06)",
      "Chat em tempo real (WebSockets)",
      "Imagens em cloud storage",
      "Notificações de negociação"],
     accent=PRIMARY_BLUE, body_size=13)
card(s, Inches(8.84), Inches(2.2), Inches(3.9), Inches(4.4), "Polimento",
     ["Refino de responsividade mobile",
      "Mais cobertura de testes",
      "Painel admin de moderação",
      "Acessibilidade"],
     accent=SUCCESS, body_size=13)
footer(s, 11)

# ============================================================
# SLIDE 12 — ENCERRAMENTO / PERGUNTAS
# ============================================================
s = add_slide(); set_bg(s)
rect(s, Inches(9.7), Inches(3.0), Inches(5), Inches(5), SECONDARY_DARK)
accent_bar(s)
tb, tf = textbox(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(2.6))
p = para(tf, first=True)
add_run(p, "Obrigado!", 48, TEXT_WHITE, bold=True)
p2 = para(tf); p2.space_before = Pt(10)
add_run(p2, "Perguntas?", 28, ACCENT_CYAN, bold=True)
p3 = para(tf); p3.space_before = Pt(18)
add_run(p3, "Tech Hub  ·  Marketplace de Tecnologia em Django", 15, TEXT_GRAY)
p4 = para(tf); p4.space_before = Pt(6)
add_run(p4, "Demonstração ao vivo disponível.", 14, TEXT_MUTED)
footer(s, 12)

# ============================================================
out = Path(__file__).resolve().parent / "Tech_Hub_Apresentacao.pptx"
prs.save(out)
print(f"OK -> {out}")
print(f"Slides gerados: {len(prs.slides._sldIdLst)}")
